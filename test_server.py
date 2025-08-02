#!/usr/bin/env python3
"""
Test script for the PPT MCP Server
"""

import asyncio
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from ppt_mcp_server import PPTExtractor, WordStructureAnalyzer, DocumentCombiner, PPTGenerator

async def create_test_files():
    """Create test PowerPoint and Word files"""
    
    # Create test PowerPoint
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Sample content for testing"
    
    # Content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Project Overview"
    
    tf = body.text_frame
    tf.text = "Main objective"
    p = tf.add_paragraph()
    p.text = "Secondary goal"
    p.level = 1
    
    ppt_path = os.path.join(tempfile.gettempdir(), "test_presentation.pptx")
    prs.save(ppt_path)
    
    # Create test Word document
    doc = Document()
    
    doc.add_heading('Project Structure Document', 0)
    doc.add_heading('1. Introduction', level=1)
    doc.add_paragraph('This document outlines the project structure and requirements.')
    doc.add_heading('2. Project Overview', level=1)
    doc.add_paragraph('The main objective is to create a comprehensive solution.')
    doc.add_heading('2.1 Goals', level=2)
    doc.add_paragraph('Primary goals include efficiency and scalability.')
    doc.add_heading('3. Implementation Plan', level=1)
    doc.add_paragraph('The implementation will follow industry best practices.')
    
    word_path = os.path.join(tempfile.gettempdir(), "test_structure.docx")
    doc.save(word_path)
    
    return ppt_path, word_path

async def test_individual_functions():
    """Test individual functions"""
    print("Creating test files...")
    ppt_path, word_path = await create_test_files()
    
    print(f"Created test files:\n  PPT: {ppt_path}\n  Word: {word_path}")
    
    # Test PPT extraction
    print("\n=== Testing PPT Extraction ===")
    try:
        extractor = PPTExtractor()
        ppt_data = extractor.extract_ppt_content(ppt_path)
        print("✓ PPT extraction successful")
        print(f"Title: {ppt_data['title']}")
        print(f"Slides: {ppt_data['total_slides']}")
        for slide in ppt_data['slides']:
            print(f"  Slide {slide['slide_number']}: {slide['title']}")
    except Exception as e:
        print(f"✗ PPT extraction failed: {e}")
        return
    
    # Test Word analysis
    print("\n=== Testing Word Structure Analysis ===")
    try:
        analyzer = WordStructureAnalyzer()
        word_data = analyzer.analyze_word_structure(word_path)
        print("✓ Word analysis successful")
        print(f"Document title: {word_data['document_title']}")
        print(f"Sections: {len(word_data['sections'])}")
        for section in word_data['sections']:
            print(f"  Section: {section['title']} (level {section['level']})")
    except Exception as e:
        print(f"✗ Word analysis failed: {e}")
        return
    
    # Test combination
    print("\n=== Testing Document Combination ===")
    try:
        combiner = DocumentCombiner()
        combined_data = combiner.combine_documents(ppt_data, word_data)
        print("✓ Document combination successful")
        print(f"Slides to create: {len(combined_data['slides'])}")
        
        # Test presentation generation
        print("\n=== Testing Presentation Generation ===")
        output_path = os.path.join(tempfile.gettempdir(), "combined_presentation.pptx")
        generator = PPTGenerator()
        final_path = generator.generate_presentation(combined_data, output_path)
        
        if os.path.exists(final_path):
            print(f"✓ Combined presentation created successfully at: {final_path}")
        else:
            print("✗ Combined presentation file not found!")
            
    except Exception as e:
        print(f"✗ Document combination/generation failed: {e}")
    
    # Cleanup
    try:
        os.remove(ppt_path)
        os.remove(word_path)
        print("\n✓ Test files cleaned up")
    except:
        pass

if __name__ == "__main__":
    asyncio.run(test_individual_functions())