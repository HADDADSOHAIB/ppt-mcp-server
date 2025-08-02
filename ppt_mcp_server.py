#!/usr/bin/env python3
"""
MCP Server for PowerPoint processing and document combination.
This server can:
1. Extract information from PowerPoint files
2. Analyze Word document structures
3. Combine information intelligently based on Word document templates
4. Generate new PowerPoint presentations
"""

import asyncio
import json
import logging
import os
import tempfile
from pathlib import Path
from typing import Dict, List, Any, Optional

from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches

# Configure logging to stderr (not stdout for MCP compatibility)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Initialize FastMCP server
mcp = FastMCP("PPT Document Processor")

class PPTExtractor:
    """Extract information from PowerPoint files"""
    
    @staticmethod
    def extract_ppt_content(file_path: str) -> Dict[str, Any]:
        """Extract all content from a PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            
            presentation_data = {
                "title": "",
                "slides": [],
                "total_slides": len(prs.slides),
                "slide_layouts": []
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "title": "",
                    "content": [],
                    "images": [],
                    "tables": [],
                    "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else f"Layout {slide_idx + 1}"
                }
                
                # Extract text from all text boxes and shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            if slide_idx == 0 and not presentation_data["title"]:
                                presentation_data["title"] = text
                            
                            if shape.name.startswith("Title") or "title" in shape.name.lower():
                                slide_data["title"] = text
                            else:
                                slide_data["content"].append({
                                    "type": "text",
                                    "content": text,
                                    "shape_name": shape.name
                                })
                    
                    # Handle tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        slide_data["tables"].append(table_data)
                    
                    # Handle images
                    if shape.shape_type == 13:  # Picture type
                        slide_data["images"].append({
                            "shape_name": shape.name,
                            "size": {
                                "width": shape.width,
                                "height": shape.height
                            }
                        })
                
                presentation_data["slides"].append(slide_data)
            
            return presentation_data
            
        except Exception as e:
            logger.error(f"Error extracting PPT content: {e}")
            raise


class WordStructureAnalyzer:
    """Analyze Word document structure and extract template information"""
    
    @staticmethod
    def analyze_word_structure(file_path: str) -> Dict[str, Any]:
        """Analyze Word document structure to understand the template"""
        try:
            doc = Document(file_path)
            
            structure = {
                "document_title": "",
                "sections": [],
                "styles_used": set(),
                "structure_outline": [],
                "paragraphs_count": len(doc.paragraphs),
                "tables_count": len(doc.tables)
            }
            
            current_section = None
            section_level = 0
            
            for para_idx, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue
                
                style_name = paragraph.style.name if paragraph.style else "Normal"
                structure["styles_used"].add(style_name)
                
                # Detect headings and structure
                if "Heading" in style_name or text.isupper() or text.endswith(":"):
                    level = 1
                    if "Heading" in style_name:
                        try:
                            level = int(style_name.split()[-1])
                        except:
                            level = 1
                    
                    if level == 1 or current_section is None:
                        current_section = {
                            "title": text,
                            "level": level,
                            "content": [],
                            "subsections": []
                        }
                        structure["sections"].append(current_section)
                        section_level = level
                    else:
                        # This is a subsection
                        subsection = {
                            "title": text,
                            "level": level,
                            "content": []
                        }
                        if current_section:
                            current_section["subsections"].append(subsection)
                    
                    structure["structure_outline"].append({
                        "level": level,
                        "title": text,
                        "paragraph_index": para_idx
                    })
                else:
                    # Regular content
                    if current_section:
                        current_section["content"].append({
                            "text": text,
                            "style": style_name
                        })
                    elif para_idx == 0:
                        structure["document_title"] = text
            
            # Analyze tables
            table_structures = []
            for table_idx, table in enumerate(doc.tables):
                table_info = {
                    "table_number": table_idx + 1,
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "headers": []
                }
                
                if table.rows:
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    table_info["headers"] = headers
                
                table_structures.append(table_info)
            
            structure["table_structures"] = table_structures
            structure["styles_used"] = list(structure["styles_used"])
            
            return structure
            
        except Exception as e:
            logger.error(f"Error analyzing Word structure: {e}")
            raise


class DocumentCombiner:
    """Combine PPT content with Word structure intelligently"""
    
    @staticmethod
    def combine_documents(ppt_data: Dict[str, Any], word_structure: Dict[str, Any]) -> Dict[str, Any]:
        """Combine PowerPoint content with Word document structure"""
        
        combined_structure = {
            "presentation_title": ppt_data.get("title", word_structure.get("document_title", "Combined Presentation")),
            "slides": [],
            "source_info": {
                "ppt_slides": ppt_data.get("total_slides", 0),
                "word_sections": len(word_structure.get("sections", []))
            }
        }
        
        # Create title slide
        title_slide = {
            "slide_type": "title",
            "title": combined_structure["presentation_title"],
            "content": ["Generated from PowerPoint extraction and Word document structure"],
            "layout": "Title Slide"
        }
        combined_structure["slides"].append(title_slide)
        
        # Map Word sections to slides
        for section in word_structure.get("sections", []):
            slide = {
                "slide_type": "content",
                "title": section["title"],
                "content": [],
                "layout": "Title and Content"
            }
            
            # Add section content
            for content_item in section.get("content", []):
                slide["content"].append(content_item["text"])
            
            # Add related PPT content if available
            for ppt_slide in ppt_data.get("slides", []):
                # Try to match content based on keywords or similarity
                if DocumentCombiner._content_similarity(section["title"], ppt_slide.get("title", "")):
                    slide["content"].extend([item["content"] for item in ppt_slide.get("content", []) if item["type"] == "text"])
            
            # Add subsections as bullet points
            for subsection in section.get("subsections", []):
                slide["content"].append(f"• {subsection['title']}")
                for sub_content in subsection.get("content", []):
                    slide["content"].append(f"  - {sub_content['text']}")
            
            combined_structure["slides"].append(slide)
        
        # Add any remaining PPT slides that weren't matched
        for ppt_slide in ppt_data.get("slides", []):
            if not any(DocumentCombiner._content_similarity(ppt_slide.get("title", ""), slide["title"]) 
                      for slide in combined_structure["slides"] if slide["slide_type"] == "content"):
                slide = {
                    "slide_type": "content",
                    "title": ppt_slide.get("title", f"Slide {ppt_slide['slide_number']}"),
                    "content": [item["content"] for item in ppt_slide.get("content", []) if item["type"] == "text"],
                    "layout": "Title and Content",
                    "source": "original_ppt"
                }
                combined_structure["slides"].append(slide)
        
        return combined_structure
    
    @staticmethod
    def _content_similarity(text1: str, text2: str) -> bool:
        """Simple content similarity check"""
        if not text1 or not text2:
            return False
        
        words1 = set(text1.lower().split())
        words2 = set(text2.lower().split())
        
        if len(words1) == 0 or len(words2) == 0:
            return False
        
        intersection = words1.intersection(words2)
        return len(intersection) > 0


class PPTGenerator:
    """Generate PowerPoint presentations from combined data"""
    
    @staticmethod
    def generate_presentation(combined_data: Dict[str, Any], output_path: str) -> str:
        """Generate a new PowerPoint presentation"""
        try:
            prs = Presentation()
            
            for slide_data in combined_data.get("slides", []):
                if slide_data["slide_type"] == "title":
                    slide_layout = prs.slide_layouts[0]  # Title slide layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    title.text = slide_data["title"]
                    
                    if slide.shapes.placeholders and len(slide.shapes.placeholders) > 1:
                        subtitle = slide.shapes.placeholders[1]
                        subtitle.text = "\n".join(slide_data.get("content", []))
                
                else:  # Content slide
                    slide_layout = prs.slide_layouts[1]  # Title and content layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    title.text = slide_data["title"]
                    
                    content_placeholder = slide.shapes.placeholders[1]
                    content_text = "\n".join([f"• {item}" if not item.startswith("•") and not item.startswith("  -") 
                                            else item for item in slide_data.get("content", [])])
                    content_placeholder.text = content_text
            
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise


@mcp.tool()
async def extract_ppt_info(file_path: str) -> Dict[str, Any]:
    """
    Extract information from a PowerPoint file.
    
    Args:
        file_path: Path to the PowerPoint file (.pptx)
    
    Returns:
        Dictionary containing extracted presentation data
    """
    try:
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        if not file_path.lower().endswith(('.pptx', '.ppt')):
            return {"error": "File must be a PowerPoint presentation (.pptx or .ppt)"}
        
        extractor = PPTExtractor()
        result = extractor.extract_ppt_content(file_path)
        
        return {
            "success": True,
            "data": result,
            "message": f"Successfully extracted content from {result['total_slides']} slides"
        }
    
    except Exception as e:
        logger.error(f"Error in extract_ppt_info: {e}")
        return {"error": str(e)}


@mcp.tool()
async def analyze_word_structure(file_path: str) -> Dict[str, Any]:
    """
    Analyze the structure of a Word document to understand its template.
    
    Args:
        file_path: Path to the Word document (.docx)
    
    Returns:
        Dictionary containing document structure analysis
    """
    try:
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        if not file_path.lower().endswith(('.docx', '.doc')):
            return {"error": "File must be a Word document (.docx or .doc)"}
        
        analyzer = WordStructureAnalyzer()
        result = analyzer.analyze_word_structure(file_path)
        
        return {
            "success": True,
            "data": result,
            "message": f"Successfully analyzed structure with {len(result['sections'])} sections"
        }
    
    except Exception as e:
        logger.error(f"Error in analyze_word_structure: {e}")
        return {"error": str(e)}


@mcp.tool()
async def combine_ppt_and_word(ppt_file: str, word_file: str, output_file: Optional[str] = None) -> Dict[str, Any]:
    """
    Combine PowerPoint content with Word document structure to create a new presentation.
    
    Args:
        ppt_file: Path to the PowerPoint file (.pptx)
        word_file: Path to the Word document (.docx)
        output_file: Optional path for the output presentation (defaults to temp file)
    
    Returns:
        Dictionary with the combined presentation path and details
    """
    try:
        # Extract PPT content
        ppt_result = await extract_ppt_info(ppt_file)
        if not ppt_result.get("success"):
            return ppt_result
        
        # Analyze Word structure
        word_result = await analyze_word_structure(word_file)
        if not word_result.get("success"):
            return word_result
        
        # Combine the data
        combiner = DocumentCombiner()
        combined_data = combiner.combine_documents(ppt_result["data"], word_result["data"])
        
        # Generate output file path
        if not output_file:
            output_file = os.path.join(tempfile.gettempdir(), "combined_presentation.pptx")
        
        # Generate the presentation
        generator = PPTGenerator()
        output_path = generator.generate_presentation(combined_data, output_file)
        
        return {
            "success": True,
            "output_file": output_path,
            "slides_created": len(combined_data["slides"]),
            "combined_data": combined_data,
            "message": f"Successfully created combined presentation with {len(combined_data['slides'])} slides"
        }
    
    except Exception as e:
        logger.error(f"Error in combine_ppt_and_word: {e}")
        return {"error": str(e)}


@mcp.tool()
async def create_presentation_from_structure(structure_data: Dict[str, Any], output_file: Optional[str] = None) -> Dict[str, Any]:
    """
    Create a PowerPoint presentation from structured data.
    
    Args:
        structure_data: Dictionary containing slide structure and content
        output_file: Optional path for the output presentation
    
    Returns:
        Dictionary with the created presentation details
    """
    try:
        if not output_file:
            output_file = os.path.join(tempfile.gettempdir(), "generated_presentation.pptx")
        
        generator = PPTGenerator()
        output_path = generator.generate_presentation(structure_data, output_file)
        
        return {
            "success": True,
            "output_file": output_path,
            "slides_created": len(structure_data.get("slides", [])),
            "message": f"Successfully created presentation at {output_path}"
        }
    
    except Exception as e:
        logger.error(f"Error in create_presentation_from_structure: {e}")
        return {"error": str(e)}


if __name__ == "__main__":
    mcp.run()